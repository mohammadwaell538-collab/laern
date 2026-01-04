import os
import json
import base64
from typing import Optional, List, Dict, Any
import traceback
import io
import uuid
import time

from fastapi import FastAPI, Form, File, UploadFile, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from starlette.responses import Response

# Use google-genai (new client library)
import google.genai


app = FastAPI()

# CORS: allow all origins/methods/headers for browser usage
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure google-genai client with API key from environment variables only.
# IMPORTANT: do NOT embed API keys in source code. Set the environment variable GOOGLE_API_KEY or GENAI_API_KEY in your deployment.
GENAI_API_KEY = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GENAI_API_KEY")

# (Optional) Support reading from a local .env file during development if python-dotenv is installed.
try:
    from dotenv import load_dotenv
    load_dotenv()
    # Re-read after loading .env
    if not GENAI_API_KEY:
        GENAI_API_KEY = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GENAI_API_KEY")
except Exception:
    # python-dotenv not installed; skip
    pass

# Helper to print safely even on consoles with limited encodings
def safe_print(msg: str, *args, **kwargs):
    try:
        print(msg, *args, **kwargs)
    except UnicodeEncodeError:
        # Fallback: replace unencodable characters so printing never crashes the process
        safe = msg.encode('utf-8', errors='replace').decode('ascii', errors='replace')
        print(safe, *args, **kwargs)


def make_genai_client():
    """Create and return a google-genai client or None if not configured."""
    if not GENAI_API_KEY:
        safe_print("WARNING: GOOGLE_API_KEY/GENAI_API_KEY not found; running without Gemini.")
        return None
    try:
        client = google.genai.Client(api_key=GENAI_API_KEY)
        return client
    except Exception as e:
        safe_print(f"Error initializing google-genai client: {e}")
        return None

# Create client from loader (do not hard-code keys here in production)
client = make_genai_client()

def call_gemini(prompt: str) -> tuple[str, int, int]:
    """
    Call Google Gemini API to generate analysis.
    Returns a tuple: (analysis_text, main_topics_count, key_points_count)
    """
    if client is None:
        return ("خطأ: لم يتم تهيئة عميل Gemini (مفتاح غير متوفر)", 0, 0)

    try:
        # Use the google-genai client library to generate content (defensive)
        try:
            response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=prompt,
            )
        except AttributeError:
            # Some client versions use different API surface
            response = client.generate(prompt)

        # Extract text from response
        analysis_text = ""
        if hasattr(response, 'text') and response.text:
            analysis_text = response.text
        elif isinstance(response, dict):
            candidates = response.get('candidates') or response.get('output')
            if candidates and isinstance(candidates, list) and len(candidates) > 0:
                first = candidates[0]
                if isinstance(first, dict):
                    analysis_text = first.get('content') or first.get('text') or str(first)
        else:
            analysis_text = str(response)

        # Count approximate topics and key points (simple heuristic)
        lines = analysis_text.split('\n')
        main_topics_count = sum(1 for line in lines if line.strip().startswith(('•', '-', '●', '○')))
        key_points_count = len(lines) // 3 if main_topics_count == 0 else main_topics_count

        return analysis_text, max(5, main_topics_count), max(12, key_points_count)

    except Exception as e:
        error_msg = f"Error connecting to Gemini: {str(e)}"
        safe_print(f"Gemini Error: {error_msg}")
        return f"Error: {str(e)}", 0, 0


# Local fallback analyzer to avoid returning errors to the browser
import re

def local_fallback_analyze(text: str, youtube: str, file_content: str, analysis_type: str) -> tuple[str, int, int, str]:
    """Simple local analyzer used when Gemini is unavailable.

    Returns: (analysis_text, main_topics_count, key_points_count, note)
    """
    content = (text or "").strip()
    if not content:
        # No usable content at all
        return ("لا توجد بيانات كافية لتنفيذ تحليل محلي.", 0, 0, "لا تتوفر بيانات كافية")

    # Split into sentences and make a short summary (up to 3 sentences)
    sentences = re.split(r'(?<=[\.؟!])\s+', content)
    if len(sentences) == 1:
        # Short single-sentence content: use it as-is
        summary = sentences[0]
    else:
        summary = ' '.join(sentences[:3])

    # Heuristic counts based on sentence count (bounded)
    main_topics_count = max(1, min(5, len(sentences)))
    key_points_count = main_topics_count * 3

    # Provide a clearer note to the UI; mention if content is very short
    if len(content) < 50:
        note = "تم استخدام محلل محلي بديل (المحتوى قصير جداً، قد تكون النتائج محدودة)."
    else:
        note = "تم استخدام محلل محلي بديل لأن خدمة Gemini غير متاحة أو تجاوزت الحصة."

    return summary, main_topics_count, key_points_count, note


def local_generate_questions(text: str, max_q: int = 20, source_label: str = 'local') -> List[Dict[str, Any]]:
    """Create simple interactive questions from text (used when Gemini is not available or for quick questions).

    Returns a list of question dicts: {question, options, answer_index, source}
    """
    s = (text or '').strip()
    if not s:
        return []
    # Very naive question generation: split into sentences and create questions from each.
    sentences = re.split(r'(?<=[\.؟!])\s+', s)
    qlist: List[Dict[str, Any]] = []
    count = 0
    for i in range(len(sentences)):
        if count >= max_q:
            break
        sent = sentences[i].strip()
        if not sent:
            continue
        # create different question templates by rotating
        if i % 3 == 0:
            qtext = f"ما الفكرة الأساسية من العبارة: '{sent[:80]}'?"
            options = ["الفكرة الأساسية", "تفصيل ثانوي", "مثال توضيحي", "لا علاقة"]
            answer_index = 0
        elif i % 3 == 1:
            qtext = f"اختر الجواب الصحيح للسؤال التالي المستمد من النص: '{sent[:80]}'"
            options = ["إجابة صحيحة", "خيار خاطئ 1", "خيار خاطئ 2", "خيار خاطئ 3"]
            answer_index = 0
        else:
            qtext = f"ضع شرحًا مختصرًا للعبارة: '{sent[:80]}'"
            options = []
            answer_index = -1

        qlist.append({'question': qtext, 'options': options, 'answer_index': answer_index, 'source': source_label})
        count += 1

    # If not enough sentences to reach max_q, create synthetic questions
    while count < max_q:
        qlist.append({'question': f'سؤال نمطي رقم {count+1}', 'options': ["نعم", "لا"], 'answer_index': 0, 'source': source_label})
        count += 1

    return qlist

# --- File parsers and async job system (DOCX/PPTX/Images/Audio + job queue) ---
try:
    from docx import Document  # python-docx
except Exception:
    Document = None

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

# Image OCR
try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

# Audio transcription
try:
    from pydub import AudioSegment
except Exception:
    AudioSegment = None

try:
    import whisper
except Exception:
    whisper = None

# In-memory job store: {job_id: {status, progress, result, created_at}}
JOBS: Dict[str, Dict[str, Any]] = {}


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a .docx file bytes."""
    if Document is None:
        return ""
    try:
        bio = io.BytesIO(file_bytes)
        doc = Document(bio)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        return '\n'.join(paragraphs)
    except Exception:
        return ""


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a .pptx file bytes (slides shapes)."""
    if Presentation is None:
        return ""
    try:
        bio = io.BytesIO(file_bytes)
        prs = Presentation(bio)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text:
                        texts.append(shape.text.strip())
                except Exception:
                    continue
        return '\n'.join(t for t in texts if t)
    except Exception:
        return ""


def extract_text_from_image(file_bytes: bytes) -> str:
    """Extract text from common image formats using pytesseract (if available).

    This is defensive: if pytesseract or Tesseract binary is missing, return empty string.
    """
    if Image is None or pytesseract is None:
        return ""
    try:
        bio = io.BytesIO(file_bytes)
        img = Image.open(bio).convert('RGB')
        # Determine languages support defensively
        lang_arg = None
        try:
            # get_languages may not be available on all installs; guard it
            langs = []
            try:
                langs = pytesseract.get_languages(config='') or []
            except Exception:
                # older versions or missing Tesseract binary
                langs = []
            if any(l for l in langs if 'ara' in l or l == 'ara'):
                lang_arg = 'ara+eng'
        except Exception:
            lang_arg = None

        if lang_arg:
            text = pytesseract.image_to_string(img, lang=lang_arg)
        else:
            text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception:
        return ""


def transcribe_audio_bytes(file_bytes: bytes) -> str:
    """Transcribe audio bytes using whisper if available. Returns empty string on failure or if not available."""
    if whisper is None:
        return ""
    try:
        # Save to a temp file and let whisper handle format detection
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.wav', delete=True) as tf:
            tf.write(file_bytes)
            tf.flush()
            model = whisper.load_model('base')
            result = model.transcribe(tf.name)
            return result.get('text', '').strip()
    except Exception:
        return ""


def chunk_text(text: str, max_chars: int = 1500) -> List[str]:
    """Split text into chunks of at most max_chars (attempt sentence boundaries)."""
    if not text:
        return []
    text = text.strip()
    if len(text) <= max_chars:
        return [text]

    sentences = re.split(r'(?<=[\.؟!])\s+', text)
    chunks: List[str] = []
    cur = []
    cur_len = 0
    for s in sentences:
        slen = len(s) + 1
        if cur_len + slen > max_chars and cur:
            chunks.append(' '.join(cur))
            cur = [s]
            cur_len = slen
        else:
            cur.append(s)
            cur_len += slen
    if cur:
        chunks.append(' '.join(cur))
    return chunks


def _process_job(job_id: str, text: str, youtube: str, files: List[UploadFile], analysis_type: str):
    """Background worker: extract content, chunk, call Gemini or fallback, aggregate and store result."""
    JOBS[job_id] = {"status": "running", "progress": 0, "result": None, "created_at": time.time()}

    try:
        # Notes list should exist before file extraction so we can append to it during processing
        notes: List[str] = []

        # Extract files
        extracted_texts: List[str] = []
        for f in files or []:
            raw = f.file.read() if hasattr(f, 'file') else b''
            name = getattr(f, 'filename', '')
            if name.lower().endswith('.docx'):
                if Document is None:
                    notes.append('python-docx not installed; skipping .docx extraction')
                    extracted_texts.append('')
                else:
                    extracted_texts.append(extract_text_from_docx(raw))
            elif name.lower().endswith('.pptx'):
                if Presentation is None:
                    notes.append('python-pptx not installed; skipping .pptx extraction')
                    extracted_texts.append('')
                else:
                    extracted_texts.append(extract_text_from_pptx(raw))
            elif name.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
                if Image is None or pytesseract is None:
                    notes.append('OCR not available (Pillow or pytesseract missing)')
                    extracted_texts.append('')
                else:
                    extracted_texts.append(extract_text_from_image(raw))
            elif name.lower().endswith(('.mp3', '.wav', '.m4a', '.flac', '.ogg')):
                # Try to transcribe audio; if pydub needed to convert formats, attempt that
                if whisper is None and AudioSegment is None:
                    notes.append('ASR not available (whisper or pydub missing)')
                    extracted_texts.append('')
                else:
                    trans = transcribe_audio_bytes(raw)
                    if not trans and AudioSegment is not None:
                        try:
                            bio = io.BytesIO(raw)
                            seg = AudioSegment.from_file(bio)
                            out_bio = io.BytesIO()
                            seg.export(out_bio, format='wav')
                            trans = transcribe_audio_bytes(out_bio.getvalue())
                        except Exception:
                            trans = ''
                    extracted_texts.append(trans)
            else:
                # unsupported binary: skip
                continue

        combined = '\n'.join([text or ''] + extracted_texts + ([youtube] if youtube else []))
        if not combined.strip():
            JOBS[job_id]['status'] = 'failed'
            JOBS[job_id]['result'] = {'error': 'لا توجد بيانات كافية'}
            return

        chunks = chunk_text(combined, max_chars=1500)
        total = len(chunks)
        analyses: List[str] = []
        any_fallback = False

        interactive_questions: List[Dict[str, Any]] = []
        sources: List[Dict[str, Any]] = []

        for idx, chunk in enumerate(chunks, start=1):
            JOBS[job_id]['progress'] = int((idx - 1) / total * 100)
            prompt = f"{analysis_type}\n\n{chunk}"
            analysis_text, tcount, pcount = call_gemini(prompt)
            source_label = f"chunk-{idx}"

            if isinstance(analysis_text, str) and analysis_text.startswith('خطأ'):
                # fallback for this chunk
                fallback_text, ft_topics, ft_points, note = local_fallback_analyze(chunk, youtube, '', analysis_type)
                analyses.append(fallback_text)
                any_fallback = True
                notes.append(note)
                # local question generation
                qlist = local_generate_questions(chunk, max_q=20, source_label=source_label)
                interactive_questions.extend(qlist)
            else:
                analyses.append(analysis_text)
                # ask Gemini to generate questions if available (defensive)
                try:
                    q_prompt = f"أنشئ سؤالين تعليميين تفاعليين (اختياريين متعدد الاختيارات) من النص التالي مع ذكر الإجابة ومصدرها:\n\n{chunk}"
                    q_text, _, _ = call_gemini(q_prompt)
                    if isinstance(q_text, str) and not q_text.startswith('خطأ'):
                        # store raw generated text as a simple question placeholder (parsing complex structured output is beyond scope)
                        interactive_questions.append({'question': q_text, 'options': [], 'answer_index': 0, 'source': source_label})
                    else:
                        interactive_questions.extend(local_generate_questions(chunk, max_q=20, source_label=source_label))
                except Exception:
                    interactive_questions.extend(local_generate_questions(chunk, max_q=2, source_label=source_label))

            # record source snippet for this chunk
            sources.append({'source': source_label, 'snippet': chunk[:200]})

            JOBS[job_id]['progress'] = int(idx / total * 100)

        # Aggregate (simple): join pieces; optionally could re-summarize into a single refined summary
        final_analysis = '\n\n'.join(analyses)
        main_topics = max(1, sum(1 for a in analyses if a.strip()))
        key_points = main_topics * 3

        JOBS[job_id]['status'] = 'completed'
        JOBS[job_id]['result'] = {
            'analysis': final_analysis,
            'main_topics': main_topics,
            'key_points': key_points,
            'fallback': any_fallback,
            'note': '; '.join(notes) if notes else '',
            'interactive_questions': interactive_questions,
            'sources': sources,
        }
    except Exception as e:
        JOBS[job_id]['status'] = 'failed'
        JOBS[job_id]['result'] = {'error': str(e), 'trace': traceback.format_exc()}



@app.post("/analyze")
async def analyze(
    text: str = Form(''),
    youtube: str = Form(''),
    analysis_type: str = Form('summary'),
    file: Optional[UploadFile] = File(None),
):
    """
    Analyze educational content and return structured results.

    Returns JSON with:
    - analysis: main analysis text
    - main_topics: count of main topics
    - key_points: count of key points
    """

    print("--- /analyze called ---")
    try:
        # Print received fields for debugging
        print(f"نوع التحليل: {analysis_type}")
        print(f"طول النص: {len(text)} حرفاً")
        print(f"رابط YouTube: {youtube}")

        file_content = ""
        filename = None
        processing_notes: List[str] = []

        # Process file if provided (docx/pptx handled)
        if file is not None:
            filename = file.filename
            raw = await file.read()
            safe_print(f"File: {filename} ({len(raw)} bytes)")

            try:
                # try decode text files
                file_content = raw.decode("utf-8")
                safe_print("Decoded file as UTF-8")
            except Exception:
                # Attempt to extract from known formats
                if filename and filename.lower().endswith('.docx'):
                    if Document is None:
                        processing_notes.append("python-docx not installed; cannot extract DOCX content")
                        safe_print("python-docx not installed; skipping DOCX extraction")
                        file_content = ""
                    else:
                        file_content = extract_text_from_docx(raw)
                        safe_print("Extracted text from docx")
                elif filename and filename.lower().endswith('.pptx'):
                    if Presentation is None:
                        processing_notes.append("python-pptx not installed; cannot extract PPTX content")
                        safe_print("python-pptx not installed; skipping PPTX extraction")
                        file_content = ""
                    else:
                        file_content = extract_text_from_pptx(raw)
                        safe_print("Extracted text from pptx")
                elif filename and filename.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
                    if Image is None or pytesseract is None:
                        processing_notes.append("OCR not enabled: Pillow or pytesseract missing")
                        safe_print("OCR not available; Pillow or pytesseract missing")
                        file_content = ""
                    else:
                        file_content = extract_text_from_image(raw)
                        safe_print("Extracted text from image (OCR)")
                elif filename and filename.lower().endswith(('.mp3', '.wav', '.m4a', '.flac', '.ogg')):
                    if whisper is None and AudioSegment is None:
                        processing_notes.append("ASR not available: whisper or pydub missing")
                        safe_print("ASR not available; whisper or pydub missing")
                        file_content = ""
                    else:
                        # attempt transcription
                        trans = transcribe_audio_bytes(raw)
                        if not trans and AudioSegment is not None:
                            try:
                                bio = io.BytesIO(raw)
                                seg = AudioSegment.from_file(bio)
                                out_bio = io.BytesIO()
                                seg.export(out_bio, format='wav')
                                trans = transcribe_audio_bytes(out_bio.getvalue())
                            except Exception:
                                trans = ''
                        file_content = trans
                        print("✓ تمت محاولة تفريغ الصوت (ASR)")
                else:
                    b64_preview = base64.b64encode(raw[:500]).decode("ascii")
                    file_content = f"[ملف ثنائي: {filename}]"
                    print("✓ تمت معالجة الملف كملف ثنائي")


        # Build prompt based on analysis_type
        prompt_templates = {
            "concepts": "استخلص المفاهيم الأساسية من هذا المحتوى:\n",
            "questions": "أنشئ مجموعة من الأسئلة التعليمية (اختيار من متعدد وإجابات قصيرة) من هذا المحتوى:\n",
            "mindmap": "أنشئ خريطة ذهنية منظمة (هيكل هرمي) من هذا المحتوى:\n",
            "summary": "قدّم ملخصًا شاملاً وموضوعيًا لهذا المحتوى:\n",
        }

        prompt_prefix = prompt_templates.get(analysis_type.strip().lower(), prompt_templates["summary"])

        # Combine all content
        full_content = f"{text}\n\nرابط YouTube: {youtube}\n\n{file_content}".strip()
        prompt = f"{prompt_prefix}\n{full_content}"

        print(f"Prompt built (length: {len(prompt)} characters)")

        # Call Gemini and get structured response with error handling
        try:
            analysis_text, topics_count, points_count = call_gemini(prompt)
            # If call_gemini returned an error message (string starting with 'خطأ'), use a local fallback instead of returning an error
            if isinstance(analysis_text, str) and analysis_text.startswith('خطأ'):
                print(f"Gemini reported an error: {analysis_text}")
                # Use local fallback analyzer so browser receives a usable response
                fallback_text, ft_topics, ft_points, note = local_fallback_analyze(text, youtube, file_content, analysis_type)
                if processing_notes:
                    note = '; '.join([note] + processing_notes)
                payload = {
                    "analysis": f"{fallback_text}\n\n(ملاحظة: {note})",
                    "main_topics": ft_topics,
                    "key_points": ft_points,
                    "fallback": True,
                    "note": note,
                }
                return Response(json.dumps(payload, ensure_ascii=False, indent=2), media_type="application/json; charset=utf-8", status_code=200)

        except Exception as e:
            print(f"Exception while calling Gemini: {e}")
            traceback.print_exc()
            # Use local fallback on exception
            fallback_text, ft_topics, ft_points, note = local_fallback_analyze(text, youtube, file_content, analysis_type)
            if processing_notes:
                note = '; '.join([note] + processing_notes)
            payload = {
                "analysis": f"{fallback_text}\n\n(ملاحظة: {note})",
                "main_topics": ft_topics,
                "key_points": ft_points,
                "fallback": True,
                "note": note,
            }
            return Response(json.dumps(payload, ensure_ascii=False, indent=2), media_type="application/json; charset=utf-8", status_code=200)

        # Build response payload with statistics for UI cards
        payload = {
            "analysis": analysis_text,
            "main_topics": topics_count,
            "key_points": points_count,
            "fallback": False,
        }

        # Attach any processing notes (e.g., missing optional libraries)
        if processing_notes:
            payload['note'] = '; '.join(processing_notes)

        # Return UTF-8 JSON to preserve Arabic text
        return Response(
            json.dumps(payload, ensure_ascii=False, indent=2),
            media_type="application/json; charset=utf-8",
            status_code=200,
        )

    except Exception as e:
        print("Unhandled error in /analyze:")
        traceback.print_exc()
        error_payload = {
            "analysis": f"خطأ في الخادم: {str(e)}",
            "main_topics": 0,
            "key_points": 0,
        }
        return Response(
            json.dumps(error_payload, ensure_ascii=False, indent=2),
            media_type="application/json; charset=utf-8",
            status_code=500,
        )


# --- Async analyze (job-based) endpoints ---
@app.post("/analyze_async")
async def analyze_async(
    background_tasks: BackgroundTasks,
    text: str = Form(''),
    youtube: str = Form(''),
    analysis_type: str = Form('summary'),
    files: Optional[List[UploadFile]] = File(None),
):
    """Start an async job to analyze potentially large files. Returns job_id immediately."""
    try:
        job_id = str(uuid.uuid4())
        JOBS[job_id] = {"status": "pending", "progress": 0, "result": None, "created_at": time.time()}

        # Materialize files for background processing: read bytes into memory
        # We create a simple object with .filename and .file to avoid UploadFile constructor differences
        materialized_files: List[Any] = []
        if files:
            print(f"/analyze_async received {len(files)} files: {[type(f) for f in files]}")
            for idx, f in enumerate(files):
                try:
                    print(f"Reading file {idx}: filename={getattr(f, 'filename', None)} type={type(f)}")
                    content = await f.read()
                    class _SimpleFile:
                        def __init__(self, filename: str, content_bytes: bytes):
                            self.filename = filename
                            self.file = io.BytesIO(content_bytes)
                    materialized_files.append(_SimpleFile(f.filename, content))
                except Exception as e:
                    print(f"Error reading uploaded file {idx}: {e}")
                    raise

        background_tasks.add_task(_process_job, job_id, text, youtube, materialized_files, analysis_type)
        return {"job_id": job_id}

    except Exception as e:
        print("Unhandled error in /analyze_async:")
        traceback.print_exc()
        return Response(
            json.dumps({"error": str(e)}, ensure_ascii=False),
            media_type="application/json; charset=utf-8",
            status_code=500,
        )

@app.get('/job/{job_id}/status')
async def job_status(job_id: str):
    job = JOBS.get(job_id)
    if not job:
        return {"status": "not_found"}
    return {"status": job.get('status'), "progress": job.get('progress')}


@app.get('/job/{job_id}/result')
async def job_result(job_id: str):
    job = JOBS.get(job_id)
    if not job:
        return Response(json.dumps({"error": "not_found"}, ensure_ascii=False), media_type="application/json; charset=utf-8", status_code=404)
    if job.get('status') != 'completed':
        return Response(json.dumps({"status": job.get('status')}, ensure_ascii=False), media_type="application/json; charset=utf-8", status_code=202)
    return Response(json.dumps(job.get('result'), ensure_ascii=False), media_type="application/json; charset=utf-8", status_code=200)


@app.get("/health")
async def health():
    """Health check endpoint"""
    return {"status": "ok", "service": "Laern Educational Analysis API"}


@app.get('/dependencies')
async def dependencies():
    """Return availability of optional dependencies and environment configuration."""
    deps = {
        'python-docx': Document is not None,
        'python-pptx': Presentation is not None,
        'Pillow': Image is not None,
        'pytesseract_installed': False,
        'pytesseract_module': pytesseract is not None,
        'pydub': AudioSegment is not None,
        'whisper': whisper is not None,
        'google-genai_configured': bool(GENAI_API_KEY),
    }

    # Try to detect Tesseract binary via pytesseract if module exists
    if pytesseract is not None:
        try:
            ver = pytesseract.get_tesseract_version()
            deps['pytesseract_installed'] = True
            deps['tesseract_version'] = str(ver)
        except Exception:
            deps['pytesseract_installed'] = False

    return deps


if __name__ == "__main__":
    import uvicorn
    import sys

    PORT = 8080
    # Use plain ASCII-friendly messages to avoid encoding errors when stdout is redirected
    print(f"Starting Laern server on port {PORT}...")
    print(f"Gemini client configured: {'YES' if GENAI_API_KEY else 'NO'}")

    try:
        uvicorn.run("proxy_server:app", host="0.0.0.0", port=PORT, reload=False)
    except Exception as e:
        print(f"Error starting server on port {PORT}: {e}")
        # try an alternate port to avoid port conflicts
        alt_port = PORT + 1
        try:
            print(f"Attempting to start on alternate port {alt_port}...")
            uvicorn.run("proxy_server:app", host="0.0.0.0", port=alt_port, reload=False)
        except Exception as e2:
            print(f"Failed to start on alternate ports: {e2}")
            sys.exit(1)
