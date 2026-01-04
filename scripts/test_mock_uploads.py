"""Quick test script: uploads mock files to /analyze and /analyze_async endpoints
Prints success messages for each file type and verifies server doesn't crash.
"""
import io
import os
import sys
import time
import json
import requests

BASE = os.environ.get('LAERN_BASE', 'http://127.0.0.1:8080')

TESTS = [
    {'name': 'Word', 'filename': 'test.docx', 'content': None, 'mimetype': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'},
    {'name': 'PowerPoint', 'filename': 'test.pptx', 'content': None, 'mimetype': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'},
    {'name': 'Image', 'filename': 'test.png', 'content': None, 'mimetype': 'image/png'},
    {'name': 'Audio', 'filename': 'test.wav', 'content': None, 'mimetype': 'audio/wav'},
]

# large binary to test handling of big uploads (~12MB)
LARGE_TEST = {'name': 'LargeBinary', 'filename': 'large.bin', 'content': b'0' * (12 * 1024 * 1024), 'mimetype': 'application/octet-stream'}

# helper to generate simple valid-ish files when possible
try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except Exception:
    PptxPresentation = None

try:
    from PIL import Image
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False

# generate content
if DocxDocument is not None:
    doc = DocxDocument()
    doc.add_paragraph('This is a test document. مرحبا بالعالم')
    bio = io.BytesIO()
    doc.save(bio)
    TESTS[0]['content'] = bio.getvalue()
else:
    TESTS[0]['content'] = 'This is a fake docx binary content (docx not installed)\nمرحبا'.encode('utf-8')

if PptxPresentation is not None:
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    from pptx.util import Inches
    left = top = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(1))
    tf = txBox.text_frame
    tf.text = 'Slide 1: Test محتوى'
    bio = io.BytesIO()
    prs.save(bio)
    TESTS[1]['content'] = bio.getvalue()
else:
    TESTS[1]['content'] = 'Fake pptx content\nمحتوى'.encode('utf-8')

if PIL_AVAILABLE:
    from PIL import Image as PILImage
    bio = io.BytesIO()
    img = PILImage.new('RGB', (200, 100), color=(73, 109, 137))
    img.save(bio, format='PNG')
    TESTS[2]['content'] = bio.getvalue()
else:
    # minimal PNG header + IHDR chunk (not fully valid but OK to test binary handling)
    TESTS[2]['content'] = b'\x89PNG\r\n\x1a\n' + b'0' * 200

# simple WAV header with silence data (PCM 16bit mono)
TESTS[3]['content'] = None
try:
    import wave
    bio = io.BytesIO()
    with wave.open(bio, 'wb') as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(16000)
        w.writeframes(b'\x00' * 16000)  # 1 second silence
    TESTS[3]['content'] = bio.getvalue()
except Exception:
    TESTS[3]['content'] = b'RIFF' + b'0' * 500

# function to POST to /analyze

def post_analyze(test):
    url = f"{BASE}/analyze"
    files = {}
    if test.get('content') is not None:
        files['file'] = (test['filename'], test['content'], test['mimetype'])
    data = {'text': 'Test bilingual: Hello / مرحبا', 'youtube': '', 'analysis_type': 'summary'}
    try:
        resp = requests.post(url, data=data, files=files, timeout=30)
        return resp
    except Exception as e:
        print(f"ERROR posting {test['name']} to /analyze: {e}")
        return None

# function to POST to /analyze_async and poll

def post_analyze_async(test, timeout=60):
    url = f"{BASE}/analyze_async"
    files = []
    if test.get('content') is not None:
        # Use 'file' field name to match the client behavior (single file upload)
        files.append(('file', (test['filename'], test['content'], test['mimetype'])))
    data = {'text': 'Async test bilingual: Hello / مرحبا', 'youtube': '', 'analysis_type': 'summary'}
    try:
        resp = requests.post(url, data=data, files=files, timeout=30)
    except Exception as e:
        print(f"ERROR posting {test['name']} to /analyze_async: {e}")
        return None

    if resp is None or resp.status_code not in (200, 201, 202):
        print(f"Unexpected response {resp.status_code if resp else 'None'} from /analyze_async for {test['name']}")
        return None

    try:
        job_id = resp.json().get('job_id')
    except Exception:
        print(f"Invalid JSON from /analyze_async for {test['name']}: {resp.text}")
        return None

    if not job_id:
        print(f"No job_id returned for {test['name']}")
        return None

    print(f"Async job started for {test['name']}: {job_id}")
    # poll
    start = time.time()
    while time.time() - start < timeout:
        try:
            s = requests.get(f"{BASE}/job/{job_id}/status", timeout=10)
            if s.status_code == 200:
                data = s.json()
                status = data.get('status')
                progress = data.get('progress')
                print(f"  status={status} progress={progress}")
                if status in ('completed', 'done'):
                    r = requests.get(f"{BASE}/job/{job_id}/result", timeout=30)
                    if r.status_code == 200:
                        try:
                            data = r.json()
                            if 'analysis' in data:
                                print(f"  Result for {test['name']} OK; analysis_present={bool(data.get('analysis').strip())}")
                                return True
                            else:
                                print(f"  Result missing 'analysis' field for {test['name']}: {r.text[:200]}")
                                return False
                        except Exception:
                            print(f"  Invalid JSON result for {test['name']}: {r.text[:200]}")
                            return False
                    else:
                        print(f"  Job not ready, result status: {r.status_code}")
            else:
                print(f"  Status query failed: {s.status_code}")
        except Exception as e:
            print(f"  Polling error: {e}")
        time.sleep(2)
    print(f"Timeout waiting for job {job_id} for {test['name']}")
    return False


def run_tests():
    print('Starting mock upload tests against', BASE)

    # quick health check
    try:
        h = requests.get(f"{BASE}/health", timeout=5)
        print('Health:', h.json())
    except Exception as e:
        print('Server health check failed:', e)
        return 1

    # test each type via sync analyze
    for t in TESTS:
        print('\n--- Testing sync analyze for', t['name'], '---')
        r = post_analyze(t)
        if r is None:
            print('ERROR: no response')
            continue
        try:
            data = r.json()
            if 'analysis' in data and isinstance(data.get('analysis'), str):
                print('SUCCESS:', t['name'], 'status', r.status_code, 'analysis_present', bool(data.get('analysis').strip()))
            else:
                print('ERROR: missing analysis field for', t['name'])
        except Exception:
            print('ERROR: invalid json or server error:', r.status_code, r.text[:200])

    # test async for each type
    for t in TESTS:
        print('\n--- Testing async analyze for', t['name'], '---')
        ok = post_analyze_async(t)
        if ok:
            print('ASYNC OK for', t['name'])
        else:
            print('ASYNC FAILED for', t['name'])

    # test large file (async)
    print('\n--- Testing large file upload (async) ---')
    ok = post_analyze_async(LARGE_TEST, timeout=120)
    if ok:
        print('Large async OK')
    else:
        print('Large async FAILED')

    print('\nMock upload tests complete.')
    return 0

if __name__ == '__main__':
    sys.exit(run_tests())
